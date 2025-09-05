from fastapi import FastAPI, UploadFile, File
from pydantic import BaseModel
from sentence_transformers import SentenceTransformer
import faiss
import numpy as np
import os
from fastapi.responses import FileResponse
from transformers import AutoTokenizer, pipeline, AutoModelForSeq2SeqLM
from fastapi.middleware.cors import CORSMiddleware

# For document parsing
from docx import Document
from pptx import Presentation
from pypdf import PdfReader

app = FastAPI(title="AI Knowledge Bot (Open Source)")

# Add CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Allow all origins
    allow_credentials=True,
    allow_methods=["*"],  # Allow all HTTP methods
    allow_headers=["*"],  # Allow all headers
)

# --------------------
# Setup
# --------------------
UPLOAD_DIR = "uploads"
STORE_DIR = "knowledge_store"
os.makedirs(UPLOAD_DIR, exist_ok=True)
os.makedirs(STORE_DIR, exist_ok=True)

# ---------------------------
# FAISS + Embeddings
# ---------------------------
# Sentence-transformer model (small, free)
embedding_model = SentenceTransformer('all-MiniLM-L6-v2')
index = faiss.IndexFlatL2(384)  # embedding dimension
stored_docs = []  # list of strings

# ---------------------------
# Hugging Face instruction-tuned model
# ---------------------------
# model_name = "databricks/dolly-v2-3b"
model_name = "google/flan-t5-small"
tokenizer = AutoTokenizer.from_pretrained(model_name)
# model = AutoModelForCausalLM.from_pretrained(model_name)
model=AutoModelForSeq2SeqLM.from_pretrained(model_name)
generator = pipeline(
    "text2text-generation",
    model=model,
    tokenizer=tokenizer,
    max_new_tokens=200,
    do_sample=True,
    device=-1 
)

# Load embedding model
model = SentenceTransformer("all-MiniLM-L6-v2")

# FAISS setup
dimension = 384  # dimension for all-MiniLM-L6-v2
index_file = os.path.join(STORE_DIR, "faiss.index")
docs_file = os.path.join(STORE_DIR, "docs.npy")

if os.path.exists(index_file) and os.path.exists(docs_file):
    index = faiss.read_index(index_file)
    stored_docs = np.load(docs_file, allow_pickle=True).tolist()
else:
    index = faiss.IndexFlatL2(dimension)
    stored_docs = []


# --------------------
# Utils
# --------------------
def extract_text(file_path: str, filename: str) -> str:
    """Extract text from txt, docx, pptx, pdf"""
    text = ""
    if filename.endswith(".txt"):
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()
    elif filename.endswith(".docx"):
        doc = Document(file_path)
        text = "\n".join([p.text for p in doc.paragraphs])
    elif filename.endswith(".pptx"):
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    elif filename.endswith(".pdf"):
        reader = PdfReader(file_path)
        for page in reader.pages:
            extracted = page.extract_text()
            if extracted:
                text += extracted + "\n"
    else:
        text = ""  # unsupported file type
    return text.strip()


# ---------------------------
# Upload endpoint
# ---------------------------

@app.post("/upload")
async def upload_file(file: UploadFile = File(...)):
    """Upload a document, extract text, create embeddings, and store in FAISS with chunking"""
    try:
        # Save file to uploads/
        file_location = os.path.join(UPLOAD_DIR, file.filename)
        with open(file_location, "wb") as f:
            f.write(await file.read())

        # Extract text
        text = extract_text(file_location, file.filename)
        if not text:
            return {"error": f"Unsupported or empty file: {file.filename}"}

        # ---- CHUNKING FUNCTION ----
        def chunk_text(text, chunk_size=300, overlap=50):
            words = text.split()
            chunks = []
            for i in range(0, len(words), chunk_size - overlap):
                chunk = " ".join(words[i:i + chunk_size])
                if chunk.strip():
                    chunks.append(chunk)
            return chunks

        # Create overlapping chunks
        chunks = chunk_text(text, chunk_size=300, overlap=50)

        # Generate embeddings for each chunk
        embeddings = model.encode(chunks)

        # Add to FAISS index
        index.add(np.array(embeddings, dtype=np.float32))

        # Store mapping of chunks to file
        for chunk in chunks:
            stored_docs.append({
                "text": chunk,
                "file_path": file_location,
                "filename": file.filename
            })

        # Persist index + docs
        faiss.write_index(index, index_file)
        np.save(docs_file, np.array(stored_docs, dtype=object))

        return {
            "filename": file.filename,
            "status": "uploaded and indexed",
            "chunks_added": len(chunks),
            "preview": chunks[:2]  # first 2 chunks
        }

    except Exception as e:
        return {"error": str(e)}


class QueryRequest(BaseModel):
    query: str

@app.post("/query/")
async def query_knowledge(req: QueryRequest):
    """Query knowledge base with FAISS + Flan-T5 (dedup + no repetition)"""
    if index.ntotal == 0 or len(stored_docs) == 0:
        return {
            "query": req.query,
            "answer": "Knowledge base is empty. Please upload documents first.",
            "document_link": None
        }

    # Embed query
    query_embedding = model.encode([req.query])
    k = 5
    D, I = index.search(np.array(query_embedding, dtype=np.float32), k)

    # Collect context (deduplicated)
    contexts, seen_texts, total_length = [], set(), 0
    doc_link = None
    max_context_chars = 2500

    for idx in I[0]:
        if idx < len(stored_docs):
            doc = stored_docs[idx]
            text = doc["text"] if isinstance(doc, dict) else str(doc)
            text = text.strip()

            if text in seen_texts:
                continue

            if total_length + len(text) <= max_context_chars:
                contexts.append(text)
                seen_texts.add(text)
                total_length += len(text)

                if not doc_link:
                    doc_link = f"http://localhost:8000/download/{idx}"

    context = " ".join(contexts)

    # Build prompt
    prompt = f"""
You are a helpful assistant. Use the following context to answer the question clearly and fully. 
Avoid repeating the same phrase multiple times. Rephrase in your own words. 
If the context is not enough, say so instead of guessing.

Context:
{context}

Question: {req.query}

Answer:
"""

    # Generate answer with anti-repetition controls
    response = generator(
        prompt,
        max_length=512,
        min_length=80,              # ensure some content
        num_return_sequences=1,
        do_sample=True,             # allow variation
        top_p=0.9,                  # nucleus sampling
        temperature=0.7,            # balanced creativity
        no_repeat_ngram_size=3,     # block loops
        repetition_penalty=1.2,     # lighter penalty
        early_stopping=True
    )

    return {
        "query": req.query,
        "answer": response[0]["generated_text"].strip(),
        "document_link": doc_link
    }


@app.get("/download/{doc_id}")
async def download_file(doc_id: int):
    if 0 <= doc_id < len(stored_docs):
        doc = stored_docs[doc_id]

        # Handle old data (string only)
        if isinstance(doc, str):
            return {"error": "This document has no file attached (old data format). Please re-upload."}

        file_path = doc["file_path"]
        if file_path and os.path.exists(file_path):
            return FileResponse(path=file_path, filename=doc["filename"])
        return {"error": "File not found on disk"}
    return {"error": "Invalid document ID"}


@app.get("/")
def home():
    return {"message": "Welcome to AI Knowledge Bot (open-source) 🚀"}
